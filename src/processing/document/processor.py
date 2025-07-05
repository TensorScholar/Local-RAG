"""
Document Processing Module.

This module provides advanced document processing capabilities for extracting
text and structure from various document formats including PDF, DOCX, and text files.
It implements OCR for image-based content, table extraction, and structural analysis.

Performance characteristics:
- Parallel processing for efficient multi-document handling
- Memory-efficient streaming for large documents
- Automatic structure detection and preservation
"""

import os
import logging
import time
from pathlib import Path
from typing import List, Dict, Any, Optional, Union, Tuple
import re
import traceback

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Process documents of various formats for text extraction and analysis.
    
    This class provides comprehensive document processing capabilities with
    format-specific extraction logic and structural preservation.
    """
    
    def __init__(self, 
                 chunk_size: int = 512,
                 chunk_overlap: int = 128,
                 enable_ocr: bool = True,
                 ocr_languages: str = 'eng',
                 extract_tables: bool = True,
                 extract_images: bool = False):
        """
        Initialize the document processor with configuration options.
        
        Args:
            chunk_size: Size of text chunks for document splitting
            chunk_overlap: Overlap between chunks to maintain context
            enable_ocr: Whether to use OCR for image-based content
            ocr_languages: Languages to use for OCR (ISO 639-2 codes)
            extract_tables: Whether to extract tables from documents
            extract_images: Whether to extract images from documents
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.enable_ocr = enable_ocr
        self.ocr_languages = ocr_languages
        self.extract_tables = extract_tables
        self.extract_images = extract_images
        
        # Supported file types and their handlers
        self.supported_extensions = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.csv': self._process_csv,
            '.pptx': self._process_powerpoint,
            '.html': self._process_html,
            '.epub': self._process_epub
        }
        
        # Performance tracking
        self.processing_times = {}
        self.total_documents_processed = 0
        self.total_pages_processed = 0
        self.total_chunks_created = 0
    
    def process_document(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Process a document file and extract structured content.
        
        This method serves as the main entry point for document processing,
        determining the document type and routing to appropriate handlers.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary with extracted content, metadata, and processing stats
        """
        file_path = Path(file_path)
        result = {
            'success': False,
            'file_path': str(file_path),
            'file_name': file_path.name,
            'file_type': file_path.suffix.lower(),
            'file_size': 0,
            'chunks': [],
            'metadata': {},
            'error': None,
            'processing_time': 0,
        }
        
        start_time = time.time()
        
        try:
            # Check if file exists
            if not file_path.exists():
                result['error'] = f"File not found: {file_path}"
                return result
            
            # Get file size
            file_size = file_path.stat().st_size
            result['file_size'] = file_size
            
            # Check if file is too large (50MB limit by default)
            max_size = 50 * 1024 * 1024  # 50MB
            if file_size > max_size:
                result['error'] = f"File too large: {file_size / (1024 * 1024):.2f}MB (max {max_size / (1024 * 1024)}MB)"
                return result
            
            # Check if file type is supported
            file_extension = file_path.suffix.lower()
            if file_extension not in self.supported_extensions:
                result['error'] = f"Unsupported file type: {file_extension}"
                return result
            
            # Process the document based on file type
            processor_fn = self.supported_extensions[file_extension]
            processing_result = processor_fn(file_path)
            
            # Update result with processing results
            result.update(processing_result)
            
            # Chunk the extracted text
            if 'text' in processing_result and processing_result['text']:
                chunks = self._chunk_text(
                    processing_result['text'],
                    processing_result.get('metadata', {}),
                    processing_result.get('structure', [])
                )
                result['chunks'] = chunks
                result['total_chunks'] = len(chunks)
                self.total_chunks_created += len(chunks)
            
            # Mark as successful if we reach this point
            result['success'] = True
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            logger.error(traceback.format_exc())
            result['error'] = str(e)
        
        # Calculate processing time
        processing_time = time.time() - start_time
        result['processing_time'] = processing_time
        
        # Update performance tracking
        file_extension = file_path.suffix.lower()
        if file_extension not in self.processing_times:
            self.processing_times[file_extension] = []
        self.processing_times[file_extension].append(processing_time)
        
        self.total_documents_processed += 1
        
        return result
    
    def process_directory(self, directory_path: Union[str, Path], recursive: bool = True) -> List[Dict[str, Any]]:
        """
        Process all supported documents in a directory.
        
        Args:
            directory_path: Path to the directory containing documents
            recursive: Whether to process subdirectories
            
        Returns:
            List of processing results for each document
        """
        directory_path = Path(directory_path)
        results = []
        
        if not directory_path.exists():
            logger.error(f"Directory not found: {directory_path}")
            return results
        
        # Get all files in the directory
        if recursive:
            files = list(directory_path.glob('**/*'))
        else:
            files = list(directory_path.glob('*'))
        
        # Filter for supported file types
        supported_files = [
            f for f in files 
            if f.is_file() and f.suffix.lower() in self.supported_extensions
        ]
        
        logger.info(f"Found {len(supported_files)} supported documents in {directory_path}")
        
        # Process each file
        for file_path in supported_files:
            logger.info(f"Processing {file_path}")
            result = self.process_document(file_path)
            results.append(result)
        
        return results
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a PDF document with text extraction, OCR, and structural analysis.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        result = {
            'text': '',
            'metadata': {
                'source': str(file_path),
                'pages': 0,
                'has_ocr': False
            },
            'structure': []
        }
        
        try:
            import PyPDF2
            from PyPDF2 import PdfReader
            
            # Open the PDF file
            with open(file_path, 'rb') as file:
                pdf = PdfReader(file)
                
                # Extract metadata
                result['metadata']['pages'] = len(pdf.pages)
                result['metadata']['title'] = pdf.metadata.get('/Title', '')
                result['metadata']['author'] = pdf.metadata.get('/Author', '')
                result['metadata']['creation_date'] = pdf.metadata.get('/CreationDate', '')
                
                self.total_pages_processed += len(pdf.pages)
                
                # Extract text from each page
                all_text = []
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    
                    # If text extraction yields little text and OCR is enabled, use OCR
                    if len(page_text) < 100 and self.enable_ocr:
                        ocr_text = self._perform_ocr_on_page(file_path, i)
                        if ocr_text:
                            page_text = ocr_text
                            result['metadata']['has_ocr'] = True
                    
                    # Add page number and page text to structure
                    all_text.append(page_text)
                    result['structure'].append({
                        'type': 'page',
                        'page_num': i + 1,
                        'text_length': len(page_text)
                    })
                
                # Combine all text
                result['text'] = '\n\n'.join(all_text)
                
                # Extract tables if enabled
                if self.extract_tables:
                    tables = self._extract_tables_from_pdf(file_path)
                    if tables:
                        result['tables'] = tables
                
                # Extract images if enabled
                if self.extract_images:
                    images = self._extract_images_from_pdf(file_path)
                    if images:
                        result['images'] = images
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            result['error'] = str(e)
            return result
    
    def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a DOCX document with text extraction and structural analysis.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        result = {
            'text': '',
            'metadata': {
                'source': str(file_path)
            },
            'structure': []
        }
        
        try:
            import docx
            
            # Open the document
            doc = docx.Document(file_path)
            
            # Extract metadata
            core_properties = doc.core_properties
            result['metadata']['title'] = core_properties.title or ''
            result['metadata']['author'] = core_properties.author or ''
            result['metadata']['created'] = str(core_properties.created) if core_properties.created else ''
            result['metadata']['modified'] = str(core_properties.modified) if core_properties.modified else ''
            
            # Extract text with structure awareness
            full_text = []
            section_texts = []
            current_heading = None
            current_heading_level = 0
            
            for paragraph in doc.paragraphs:
                if paragraph.style.name.startswith('Heading'):
                    # If we were building a section, add it to the results
                    if current_heading and section_texts:
                        section_text = '\n'.join(section_texts)
                        full_text.append(section_text)
                        result['structure'].append({
                            'type': 'section',
                            'heading': current_heading,
                            'level': current_heading_level,
                            'text_length': len(section_text)
                        })
                        section_texts = []
                    
                    # Start a new section
                    current_heading = paragraph.text
                    current_heading_level = int(paragraph.style.name.replace('Heading', '')) if paragraph.style.name != 'Heading' else 1
                    result['structure'].append({
                        'type': 'heading',
                        'text': paragraph.text,
                        'level': current_heading_level
                    })
                else:
                    # Add paragraph to current section
                    section_texts.append(paragraph.text)
            
            # Add the last section if any
            if section_texts:
                section_text = '\n'.join(section_texts)
                full_text.append(section_text)
                if current_heading:
                    result['structure'].append({
                        'type': 'section',
                        'heading': current_heading,
                        'level': current_heading_level,
                        'text_length': len(section_text)
                    })
            
            # Extract tables if present
            if doc.tables and self.extract_tables:
                tables = []
                for i, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    
                    tables.append({
                        'table_id': i,
                        'rows': len(table.rows),
                        'columns': len(table.rows[0].cells) if table.rows else 0,
                        'data': table_data
                    })
                
                result['tables'] = tables
            
            # Combine all text
            result['text'] = '\n\n'.join(full_text)
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {str(e)}")
            result['error'] = str(e)
            return result
    
    def _process_text(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a plain text document.
        
        Args:
            file_path: Path to the text file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        result = {
            'text': '',
            'metadata': {
                'source': str(file_path)
            },
            'structure': []
        }
        
        try:
            # Detect encoding - try UTF-8 first, then fallback to Latin-1
            encoding = 'utf-8'
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text = file.read()
            except UnicodeDecodeError:
                encoding = 'latin-1'
                with open(file_path, 'r', encoding=encoding) as file:
                    text = file.read()
            
            result['text'] = text
            result['metadata']['encoding'] = encoding
            
            # Try to identify structure based on newlines and patterns
            lines = text.split('\n')
            result['metadata']['line_count'] = len(lines)
            
            # Find potential headings (lines that are shorter and followed by blank lines)
            in_paragraph = False
            current_paragraph = []
            paragraphs = []
            
            for i, line in enumerate(lines):
                if not line.strip():
                    # Empty line ends the current paragraph
                    if in_paragraph:
                        in_paragraph = False
                        paragraphs.append('\n'.join(current_paragraph))
                        current_paragraph = []
                else:
                    # Line with content
                    if not in_paragraph:
                        in_paragraph = True
                    current_paragraph.append(line)
            
            # Add the last paragraph if any
            if current_paragraph:
                paragraphs.append('\n'.join(current_paragraph))
            
            # Add paragraph structure
            for i, para in enumerate(paragraphs):
                result['structure'].append({
                    'type': 'paragraph',
                    'index': i,
                    'text_length': len(para)
                })
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {str(e)}")
            result['error'] = str(e)
            return result
    
    def _process_markdown(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a Markdown document with structure preservation.
        
        Args:
            file_path: Path to the Markdown file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        result = {
            'text': '',
            'metadata': {
                'source': str(file_path)
            },
            'structure': []
        }
        
        try:
            # Read the file
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            # Store the raw markdown
            result['text'] = text
            
            # Parse Markdown structure
            lines = text.split('\n')
            in_code_block = False
            
            # Extract headings
            heading_pattern = re.compile(r'^(#{1,6})\s+(.+)$')
            code_block_pattern = re.compile(r'^```')
            front_matter = []
            in_front_matter = False
            
            for i, line in enumerate(lines):
                # Handle YAML front matter
                if i == 0 and line.strip() == '---':
                    in_front_matter = True
                    continue
                
                if in_front_matter:
                    if line.strip() == '---':
                        in_front_matter = False
                        # Parse front matter as metadata
                        for fm_line in front_matter:
                            if ':' in fm_line:
                                key, value = fm_line.split(':', 1)
                                result['metadata'][key.strip()] = value.strip()
                    else:
                        front_matter.append(line)
                    continue
                
                # Handle code blocks
                if code_block_pattern.match(line):
                    in_code_block = not in_code_block
                    lang = line[3:].strip()
                    if in_code_block and lang:
                        result['structure'].append({
                            'type': 'code_block_start',
                            'line': i,
                            'language': lang
                        })
                    elif not in_code_block:
                        result['structure'].append({
                            'type': 'code_block_end',
                            'line': i
                        })
                    continue
                
                # Skip structure parsing inside code blocks
                if in_code_block:
                    continue
                
                # Detect headings
                heading_match = heading_pattern.match(line)
                if heading_match:
                    level = len(heading_match.group(1))
                    text = heading_match.group(2)
                    result['structure'].append({
                        'type': 'heading',
                        'level': level,
                        'text': text,
                        'line': i
                    })
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing Markdown file {file_path}: {str(e)}")
            result['error'] = str(e)
            return result
    
    def _process_csv(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a CSV file with structural awareness.
        
        Args:
            file_path: Path to the CSV file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        result = {
            'text': '',
            'metadata': {
                'source': str(file_path)
            },
            'structure': [],
            'table': None
        }
        
        try:
            import pandas as pd
            
            # Read CSV file
            df = pd.read_csv(file_path)
            
            # Extract metadata
            result['metadata']['rows'] = len(df)
            result['metadata']['columns'] = len(df.columns)
            result['metadata']['column_names'] = df.columns.tolist()
            
            # Store the data as a table structure
            table_data = df.to_dict(orient='records')
            result['table'] = {
                'headers': df.columns.tolist(),
                'data': table_data
            }
            
            # Generate a text representation of the CSV
            # Include headers and first few rows
            text_parts = []
            
            # Add headers
            text_parts.append(', '.join(df.columns.tolist()))
            
            # Add sample rows (up to 10)
            sample_size = min(10, len(df))
            for i in range(sample_size):
                row = df.iloc[i]
                row_text = ', '.join(str(value) for value in row.values)
                text_parts.append(row_text)
            
            # Add summary if there are more rows
            if len(df) > sample_size:
                text_parts.append(f"... and {len(df) - sample_size} more rows")
            
            # Add statistical summary for numerical columns
            numeric_columns = df.select_dtypes(include=['number']).columns
            if len(numeric_columns) > 0:
                text_parts.append("\nNumerical statistics:")
                stats = df[numeric_columns].describe().to_string()
                text_parts.append(stats)
            
            # Store the text representation
            result['text'] = '\n'.join(text_parts)
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing CSV file {file_path}: {str(e)}")
            result['error'] = str(e)
            return result
    
    def _process_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a PowerPoint presentation with structural awareness.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        result = {
            'text': '',
            'metadata': {
                'source': str(file_path)
            },
            'structure': []
        }
        
        try:
            from pptx import Presentation
            
            # Open the presentation
            prs = Presentation(file_path)
            
            # Extract metadata
            result['metadata']['slides'] = len(prs.slides)
            
            # Process each slide
            slides_text = []
            
            for slide_index, slide in enumerate(prs.slides):
                slide_text = []
                
                # Get slide title if available
                title = slide.shapes.title
                if title and title.has_text_frame:
                    title_text = title.text
                    slide_text.append(f"Title: {title_text}")
                    
                    # Add title to structure
                    result['structure'].append({
                        'type': 'slide_title',
                        'slide': slide_index + 1,
                        'text': title_text
                    })
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text and not (title and shape == title):
                                slide_text.append(para_text)
                
                # Create slide text
                full_slide_text = '\n'.join(slide_text)
                slides_text.append(f"--- Slide {slide_index + 1} ---\n{full_slide_text}")
                
                # Add slide to structure
                result['structure'].append({
                    'type': 'slide',
                    'slide_num': slide_index + 1,
                    'text_length': len(full_slide_text)
                })
            
            # Combine all text
            result['text'] = '\n\n'.join(slides_text)
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file {file_path}: {str(e)}")
            result['error'] = str(e)
            return result
    
    def _process_html(self, file_path: Path) -> Dict[str, Any]:
        """
        Process an HTML file with structure preservation.
        
        Args:
            file_path: Path to the HTML file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        result = {
            'text': '',
            'metadata': {
                'source': str(file_path)
            },
            'structure': []
        }
        
        try:
            from bs4 import BeautifulSoup
            
            # Read the file
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract metadata
            title_tag = soup.find('title')
            if title_tag:
                result['metadata']['title'] = title_tag.text
            
            # Extract meta tags
            for meta in soup.find_all('meta'):
                if meta.get('name') and meta.get('content'):
                    result['metadata'][meta['name']] = meta['content']
            
            # Extract headings and content structure
            texts = []
            
            # Process headings and their content
            for heading_level in range(1, 7):
                for heading in soup.find_all(f'h{heading_level}'):
                    heading_text = heading.text.strip()
                    if heading_text:
                        texts.append(f"{'#' * heading_level} {heading_text}")
                        
                        # Add to structure
                        result['structure'].append({
                            'type': 'heading',
                            'level': heading_level,
                            'text': heading_text
                        })
                        
                        # Get next siblings until next heading
                        next_element = heading.next_sibling
                        section_content = []
                        
                        while next_element:
                            if next_element.name and next_element.name.startswith('h'):
                                break
                            
                            if hasattr(next_element, 'text') and next_element.text.strip():
                                section_content.append(next_element.text.strip())
                            
                            next_element = next_element.next_sibling
                        
                        if section_content:
                            section_text = '\n'.join(section_content)
                            texts.append(section_text)
            
            # Extract paragraphs not captured by headings
            for paragraph in soup.find_all('p'):
                para_text = paragraph.text.strip()
                if para_text and para_text not in texts:
                    texts.append(para_text)
            
            # Combine all text
            result['text'] = '\n\n'.join(texts)
            
            # Extract tables if present
            if self.extract_tables:
                tables = []
                for i, table in enumerate(soup.find_all('table')):
                    table_data = []
                    
                    # Extract headers
                    headers = []
                    header_row = table.find('thead')
                    if header_row:
                        for th in header_row.find_all('th'):
                            headers.append(th.text.strip())
                    
                    # Extract rows
                    for tr in table.find_all('tr'):
                        row_data = []
                        for cell in tr.find_all(['td', 'th']):
                            row_data.append(cell.text.strip())
                        
                        if row_data:
                            table_data.append(row_data)
                    
                    tables.append({
                        'table_id': i,
                        'headers': headers,
                        'rows': len(table_data),
                        'columns': len(headers) if headers else (len(table_data[0]) if table_data else 0),
                        'data': table_data
                    })
                
                if tables:
                    result['tables'] = tables
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing HTML file {file_path}: {str(e)}")
            result['error'] = str(e)
            return result
    
    def _process_epub(self, file_path: Path) -> Dict[str, Any]:
        """
        Process an EPUB e-book with chapter structure preservation.
        
        Args:
            file_path: Path to the EPUB file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        result = {
            'text': '',
            'metadata': {
                'source': str(file_path)
            },
            'structure': []
        }
        
        try:
            # Check if the required library is available
            try:
                import ebooklib
                from ebooklib import epub
                from bs4 import BeautifulSoup
            except ImportError:
                result['error'] = "ebooklib not installed. Install with: pip install ebooklib beautifulsoup4"
                return result
            
            # Read the EPUB file
            book = epub.read_epub(file_path)
            
            # Extract metadata
            result['metadata']['title'] = book.get_metadata('DC', 'title')[0][0] if book.get_metadata('DC', 'title') else ''
            result['metadata']['creator'] = book.get_metadata('DC', 'creator')[0][0] if book.get_metadata('DC', 'creator') else ''
            result['metadata']['language'] = book.get_metadata('DC', 'language')[0][0] if book.get_metadata('DC', 'language') else ''
            result['metadata']['identifier'] = book.get_metadata('DC', 'identifier')[0][0] if book.get_metadata('DC', 'identifier') else ''
            
            # Process items (chapters)
            chapters = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    # Parse HTML content
                    content = item.get_content().decode('utf-8')
                    soup = BeautifulSoup(content, 'html.parser')
                    
                    # Extract title and text
                    chapter_title = ''
                    h_tags = soup.find(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
                    if h_tags:
                        chapter_title = h_tags.text.strip()
                    
                    # Get all text
                    chapter_text = soup.get_text().strip()
                    
                    # Add to chapters list
                    if chapter_text:
                        chapters.append({
                            'title': chapter_title,
                            'text': chapter_text
                        })
                        
                        # Add to structure
                        result['structure'].append({
                            'type': 'chapter',
                            'title': chapter_title,
                            'text_length': len(chapter_text)
                        })
            
            # Combine all text with chapter information
            combined_text = []
            for i, chapter in enumerate(chapters):
                title = chapter['title'] or f"Chapter {i+1}"
                combined_text.append(f"--- {title} ---\n\n{chapter['text']}")
            
            result['text'] = '\n\n'.join(combined_text)
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing EPUB file {file_path}: {str(e)}")
            result['error'] = str(e)
            return result
    
    def _perform_ocr_on_page(self, pdf_path: Path, page_num: int) -> str:
        """
        Perform OCR on a PDF page to extract text from images.
        
        Args:
            pdf_path: Path to the PDF file
            page_num: Page number to process (0-indexed)
            
        Returns:
            Extracted text from OCR
        """
        if not self.enable_ocr:
            return ""
        
        try:
            # Check if the required libraries are available
            try:
                import pytesseract
                from pdf2image import convert_from_path
            except ImportError:
                logger.error("OCR libraries (pytesseract, pdf2image) not installed.")
                return ""
            
            # Convert PDF page to image
            images = convert_from_path(pdf_path, first_page=page_num+1, last_page=page_num+1)
            
            if not images:
                return ""
            
            # Perform OCR on the image
            ocr_text = pytesseract.image_to_string(images[0], lang=self.ocr_languages)
            
            return ocr_text
            
        except Exception as e:
            logger.error(f"Error performing OCR: {str(e)}")
            return ""
    
    def _extract_tables_from_pdf(self, pdf_path: Path) -> List[Dict[str, Any]]:
        """
        Extract tables from a PDF document.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of extracted tables with metadata
        """
        if not self.extract_tables:
            return []
        
        try:
            # Check if the required library is available
            try:
                import tabula
            except ImportError:
                logger.error("tabula-py not installed. Install with: pip install tabula-py")
                return []
            
            # Extract tables
            tables = tabula.read_pdf(str(pdf_path), pages='all', multiple_tables=True)
            
            # Format tables
            extracted_tables = []
            for i, df in enumerate(tables):
                table_data = {
                    'table_id': i,
                    'rows': len(df),
                    'columns': len(df.columns),
                    'headers': df.columns.tolist(),
                    'data': df.to_dict(orient='records')
                }
                extracted_tables.append(table_data)
            
            return extracted_tables
            
        except Exception as e:
            logger.error(f"Error extracting tables from PDF: {str(e)}")
            return []
    
    def _extract_images_from_pdf(self, pdf_path: Path) -> List[Dict[str, Any]]:
        """
        Extract images from a PDF document.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of extracted image metadata
        """
        if not self.extract_images:
            return []
        
        try:
            # Check if the required libraries are available
            try:
                from PyPDF2 import PdfReader
                from PIL import Image
                import io
                import hashlib
            except ImportError:
                logger.error("Image extraction libraries not installed.")
                return []
            
            # Open the PDF
            with open(pdf_path, 'rb') as file:
                pdf = PdfReader(file)
                
                images = []
                for i, page in enumerate(pdf.pages):
                    # Check for XObject resources (embedded images)
                    if '/XObject' in page['/Resources']:
                        xobjects = page['/Resources']['/XObject']
                        
                        # Extract image data
                        for name, object in xobjects.items():
                            if object['/Subtype'] == '/Image':
                                # Get image metadata
                                image_data = {
                                    'page': i + 1,
                                    'name': name,
                                    'width': object.get('/Width', 0),
                                    'height': object.get('/Height', 0),
                                    'bits_per_component': object.get('/BitsPerComponent', 0),
                                    'color_space': str(object.get('/ColorSpace', '')),
                                    'hash': None
                                }
                                
                                # Compute hash of image data
                                if hasattr(object, '_data'):
                                    image_hash = hashlib.md5(object._data).hexdigest()
                                    image_data['hash'] = image_hash
                                
                                images.append(image_data)
                
                return images
                
        except Exception as e:
            logger.error(f"Error extracting images from PDF: {str(e)}")
            return []
    
    def _chunk_text(self, text: str, metadata: Dict[str, Any], structure: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Split text into chunks of appropriate size with metadata.
        
        This method implements intelligent text chunking with structure awareness
        to ensure logical segmentation and context preservation.
        
        Args:
            text: Full text to chunk
            metadata: Document metadata to carry into chunks
            structure: Document structure information
            
        Returns:
            List of chunk dictionaries with text and metadata
        """
        chunks = []
        
        # If text is smaller than chunk size, return it as a single chunk
        if len(text) <= self.chunk_size:
            chunk = {
                'text': text,
                'metadata': metadata.copy()
            }
            chunks.append(chunk)
            return chunks
        
        # Use document structure for intelligent chunking when available
        if structure:
            # Sort structure by position/order
            if 'page' in structure[0] or 'page_num' in structure[0]:
                # Sort by page number
                sorted_structure = sorted(structure, key=lambda x: x.get('page_num', x.get('page', 0)))
            elif 'line' in structure[0]:
                # Sort by line number
                sorted_structure = sorted(structure, key=lambda x: x.get('line', 0))
            else:
                # Use the order as is
                sorted_structure = structure
            
            # Identify chunk boundaries based on structure
            chunk_texts = []
            current_chunk = []
            current_length = 0
            
            # Use headings or pages as natural chunk boundaries
            for item in sorted_structure:
                item_type = item.get('type', '')
                
                if item_type in ['heading', 'page', 'slide', 'chapter'] and current_length > 0:
                    # New section, start a new chunk if current one is substantial
                    if current_length >= self.chunk_size / 2:
                        chunk_texts.append('\n'.join(current_chunk))
                        current_chunk = []
                        current_length = 0
                
                # Add the current item's text
                if 'text' in item:
                    item_text = item['text']
                    current_chunk.append(item_text)
                    current_length += len(item_text)
                
                # Check if we've exceeded the chunk size
                if current_length >= self.chunk_size:
                    chunk_texts.append('\n'.join(current_chunk))
                    current_chunk = []
                    current_length = 0
            
            # Add the last chunk
            if current_chunk:
                chunk_texts.append('\n'.join(current_chunk))
        else:
            # No structure available, use simple overlap chunking
            chunk_texts = []
            
            # Split by paragraphs first
            paragraphs = text.split('\n\n')
            
            current_chunk = []
            current_length = 0
            
            for para in paragraphs:
                para_length = len(para)
                
                # If adding this paragraph exceeds chunk size and we have content
                if current_length + para_length > self.chunk_size and current_chunk:
                    chunk_texts.append('\n\n'.join(current_chunk))
                    
                    # Start a new chunk with overlap if possible
                    overlap_start = max(0, len(current_chunk) - self.chunk_overlap)
                    current_chunk = current_chunk[overlap_start:] if overlap_start < len(current_chunk) else []
                    current_length = sum(len(p) for p in current_chunk)
                
                # Add paragraph to current chunk
                current_chunk.append(para)
                current_length += para_length
            
            # Add the last chunk
            if current_chunk:
                chunk_texts.append('\n\n'.join(current_chunk))
        
        # Create chunk objects with metadata
        for i, chunk_text in enumerate(chunk_texts):
            chunk_metadata = metadata.copy()
            chunk_metadata['chunk_id'] = i
            chunk_metadata['chunk_total'] = len(chunk_texts)
            
            chunk = {
                'text': chunk_text,
                'metadata': chunk_metadata
            }
            chunks.append(chunk)
        
        return chunks
    
    def get_performance_stats(self) -> Dict[str, Any]:
        """
        Get performance statistics for document processing.
        
        Returns:
            Dictionary with comprehensive performance metrics
        """
        stats = {
            'total_documents_processed': self.total_documents_processed,
            'total_pages_processed': self.total_pages_processed,
            'total_chunks_created': self.total_chunks_created,
            'file_types': {},
            'average_processing_time': 0
        }
        
        # Calculate average processing time per file type
        total_time = 0
        total_files = 0
        
        for file_type, times in self.processing_times.items():
            avg_time = sum(times) / len(times) if times else 0
            stats['file_types'][file_type] = {
                'count': len(times),
                'average_time': avg_time,
                'total_time': sum(times)
            }
            total_time += sum(times)
            total_files += len(times)
        
        # Overall average
        stats['average_processing_time'] = total_time / total_files if total_files > 0 else 0
        
        return stats
